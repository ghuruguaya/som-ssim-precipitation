# -*- coding: utf-8 -*-
"""
Created on Sun Oct 26 08:54:47 2025

@author: Gaby
"""

#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Desnormaliza los PESOS del SOM (weights.npy) de escala Min–Max (Oct–Mar) a mm
y exporta un mapa por nodo con estética cuidada (Cartopy):

- Cartopy PlateCarree
- Gridlines con etiquetas (lon abajo, lat izquierda)
- COSTAS, BORDERS, STATES
- Colorbar horizontal
- Colormap RdBu_r (positivas azules, negativas rojas)
- Océano blanco (NaN)
- Rango clipeado [-150, 150] mm
"""

import os, json, glob
import numpy as np
import xarray as xr
import matplotlib.pyplot as plt
import cartopy.crs as ccrs
import cartopy.feature as cfeature
from matplotlib.colors import Normalize
from matplotlib.ticker import FuncFormatter
from matplotlib import colormaps as cm  # <- en lugar de get_cmap

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False

# =======================
# ======= CONFIG =========
# =======================
BASE = r"C:\Users\Gaby\Desktop\SOM\SSIM-COMPLETO\PROYECTO_GABY"

# NetCDF de anomalías en mm (si ya está recortado a Oct–Mar, perfecto)
ANOM_NC = os.path.join(BASE, "chirps_anom_oct_mar.nc")

RESULTS_DIR = os.path.join(BASE, "som_results_ssim")
CFG_NAME = "som_4x4_pca_bmu-ssim"   # si dudás, dejalo "" y autodetecta

OUT_SUBDIR = "weights_mm"

VMIN, VMAX = -150, 150
LEVELS = np.linspace(VMIN, VMAX, 21)
CMAP = cm.get_cmap("RdBu").copy(); CMAP.set_bad("white")
FIGSIZE = (6.0, 8.0)
DPI = 150
SAVE_PPTX = True  # si está python-pptx, hace PPT con todos los nodos

# =======================
# ===== UTILIDADES ======
# =======================
def autodetect_cfg_dir(results_dir, cfg_name):
    if cfg_name:
        cfg_dir = os.path.join(results_dir, cfg_name)
        if os.path.exists(os.path.join(cfg_dir, "weights.npy")):
            return cfg_dir
    cands = []
    for d in glob.glob(os.path.join(results_dir, "*")):
        if os.path.isdir(d) and os.path.exists(os.path.join(d, "weights.npy")):
            cands.append(d)
    if not cands:
        raise FileNotFoundError("No encontré subcarpetas con weights.npy dentro de som_results_ssim.")
    cands.sort(key=lambda p: os.path.getmtime(p), reverse=True)
    return cands[0]

def pick_var(ds):
    for v in ds.data_vars:
        return v
    raise ValueError("No encontré variable de datos en el NetCDF.")

def guess_coords(ds):
    lat_names = ["lat", "latitude", "y"]
    lon_names = ["lon", "longitude", "x"]
    time_names = ["time", "t"]
    latn = next((n for n in lat_names if n in ds.coords or n in ds.dims), None)
    lonn = next((n for n in lon_names if n in ds.coords or n in ds.dims), None)
    timn = next((n for n in time_names if n in ds.coords or n in ds.dims), None)
    if not (latn and lonn and timn):
        raise ValueError("No pude identificar coords lat/lon/time en el NetCDF.")
    return latn, lonn, timn

def plot_mapa(ax, LON, LAT, data2d, titulo):
    ax.set_global()
    gl = ax.gridlines(draw_labels=True, linewidth=0.5, color='gray', alpha=0.5, linestyle='--')
    gl.top_labels = False
    gl.right_labels = False
    gl.bottom_labels = True
    gl.left_labels = True
    gl.xformatter = FuncFormatter(lambda x, pos: f"{int(x)}")
    gl.yformatter = FuncFormatter(lambda y, pos: f"{int(y)}")
    gl.xlabel_style = {'size': 11}
    gl.ylabel_style = {'size': 11}

    data2d = np.ma.masked_invalid(data2d)
    im = ax.contourf(LON, LAT, data2d, levels=LEVELS, cmap=CMAP,
                     norm=Normalize(VMIN, VMAX), transform=ccrs.PlateCarree(),
                     extend='both', antialiased=True)

    ax.coastlines()
    ax.add_feature(cfeature.BORDERS, linestyle=':', linewidth=0.6)
    ax.add_feature(cfeature.STATES, linestyle=':', edgecolor='gray', linewidth=0.4)
    ax.set_extent([np.nanmin(LON), np.nanmax(LON), np.nanmin(LAT), np.nanmax(LAT)], crs=ccrs.PlateCarree())
    ax.set_title(titulo, fontsize=13)
    return im

def vec_to_2d(vec, spatial_shape, valid_idx):
    M = np.full(spatial_shape, np.nan, dtype=float)
    M[valid_idx] = vec
    return M

# =======================
# ========= MAIN ========
# =======================
def main():
    # 1) Abrir NetCDF (mm) + coords
    if not os.path.exists(ANOM_NC):
        raise FileNotFoundError(f"No encontré el NetCDF de anomalías: {ANOM_NC}")
    ds = xr.open_dataset(ANOM_NC)
    vname = pick_var(ds)
    latn, lonn, timn = guess_coords(ds)
    da = ds[vname]  # típico: (time, lat, lon)
    lats = ds[latn].values
    lons = ds[lonn].values

    # 2) Filtrar a Oct–Mar si hiciera falta (si ya está recortado, no cambia)
    if timn in ds.coords:
        t = ds[timn].dt  # acceso vectorizado a componentes de fecha
        da = da.where((t.month >= 10) | (t.month <= 3), drop=True)
    else:
        # fallback si 'time' no es coord
        tiempo = xr.DataArray(pd.to_datetime(ds[timn].values), dims=[timn])
        da = da.assign_coords({timn: tiempo})
        da = da.where((da[timn].dt.month >= 10) | (da[timn].dt.month <= 3), drop=True)

    # 3) Máscara de océano: NaN donde todo es NaN en el tiempo (Oct–Mar)
    #    Evita warnings de "all-NaN slice" en min/max
    mask_allnan = ~np.isfinite(da.max(dim=timn, skipna=True))
    # 4) Min/Max por píxel usando xarray con skipna (sin warnings)
    min_mm = da.min(dim=timn, skipna=True).where(~mask_allnan).values
    max_mm = da.max(dim=timn, skipna=True).where(~mask_allnan).values

    # 5) LON/LAT grid para plot
    if lats.ndim == 1 and lons.ndim == 1:
        LON, LAT = np.meshgrid(lons, lats)
    else:
        LON, LAT = lons, lats

    # 6) Cargar SOM: weights + valid_idx + config
    cfg_dir = autodetect_cfg_dir(RESULTS_DIR, CFG_NAME)
    cfg_name = os.path.basename(cfg_dir)

    with open(os.path.join(cfg_dir, "config.json")) as f:
        cfg = json.load(f)

    W = np.load(os.path.join(cfg_dir, "weights.npy"))              # [rows, cols, n_valid]
    vidx = np.load(os.path.join(cfg_dir, "valid_idx.npy"), allow_pickle=True)
    # ⚠️ asegurar enteros
    valid_idx = (np.asarray(vidx[0]).astype(int), np.asarray(vidx[1]).astype(int))

    rows, cols = cfg["grid_size"]
    # Usar shape espacial directamente del NetCDF
    spatial_shape = (len(lats), len(lons))

    # Sanity check: longitudes de vector vs cantidad de válidos
    n_valid = len(valid_idx[0])
    if W.shape[-1] != n_valid:
        raise ValueError(f"weights[..., n_valid]={W.shape[-1]} no coincide con len(valid_idx)={n_valid}")

    # 7) Desnormalizar: w_mm = w_norm*(max-min) + min
    span = (max_mm - min_mm)
    # donde span=0 (o NaN) dejamos NaN
    span = np.where((span == 0) | ~np.isfinite(span), np.nan, span)

    out_base = os.path.join(cfg_dir, OUT_SUBDIR)
    os.makedirs(out_base, exist_ok=True)

    for i in range(rows):
        for j in range(cols):
            w2d_norm = vec_to_2d(W[i, j], spatial_shape, valid_idx)  # [0,1]
            w2d_mm   = w2d_norm * span + min_mm                      # mm
            w2d_mm   = np.clip(w2d_mm, VMIN, VMAX)

            fig = plt.figure(figsize=FIGSIZE, dpi=DPI)
            ax = plt.axes(projection=ccrs.PlateCarree())
            im = plot_mapa(ax, LON, LAT, w2d_mm, f"Nodo ({i},{j})")
            cbar = plt.colorbar(im, orientation='horizontal', pad=0.05, extend='both')
            cbar.set_label("Anomalía (mm)", fontsize=11)
            plt.tight_layout()

            fname = os.path.join(out_base, f"weight_mm_nodo_{i}_{j}.png")
            plt.savefig(fname)
            plt.close(fig)
            print(f"✓ guardado {fname}")

    if SAVE_PPTX and HAVE_PPTX:
        ppt = Presentation()
        for i in range(rows):
            for j in range(cols):
                img = os.path.join(out_base, f"weight_mm_nodo_{i}_{j}.png")
                if not os.path.exists(img): 
                    continue
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                slide.shapes.title.text = f"Nodo ({i},{j}) — {cfg_name}"
                left, top, height = Inches(1), Inches(1.2), Inches(5.8)
                slide.shapes.add_picture(img, left, top, height=height)
        ppt_path = os.path.join(out_base, f"{cfg_name}_weights_mm.pptx")
        ppt.save(ppt_path)
        print(f"📑 PPT guardado: {ppt_path}")

    print("✅ Listo. Pesos desnormalizados exportados a:", out_base)

if __name__ == "__main__":
    # Spyder a veces necesita estos imports adentro
    import pandas as pd
    main()
